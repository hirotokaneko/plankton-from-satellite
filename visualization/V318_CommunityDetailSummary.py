import pandas as pd
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

#Arguments
SIZ = sys.argv[1]
OPO = sys.argv[2]
TPR = sys.argv[3]
FRQ = sys.argv[4]
NRM = sys.argv[5]

#Set file name variables
opo2 = "" if OPO == "all" else "{}.".format(OPO)
tpr2 = "" if TPR == "keep" else "thinned{}.".format(TPR)
subdirec = "{}.{}{}frq{}".format(SIZ,opo2,tpr2,FRQ)

#Make a presentation
prs = Presentation()
prs.slide_height = Inches(9)
prs.slide_width = Inches(16)

#Set balank slide layout
blank_slide_layout = prs.slide_layouts[6]
font = "Helvetica"

#Read the table of communities
DFcom = pd.read_csv("../data/{}/community.{}.{}{}frq{}.{}.tsv".format(subdirec,SIZ,opo2,tpr2,FRQ,NRM), delimiter="\t", index_col=0)

#Function for making a text box
def MakeTextBox(top, left, text, size):
	width = height = Inches(1)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = text
	p.font.size = size
	p.font.name = font
	return tf

#Function for adding a row to the text box
def AddTextToBox(tf, text, size):
	p = tf.add_paragraph()
	p.text = text
	p.font.size = size
	p.font.name = font

#Loop by community detection methods
list_com_detec_met = ["fast_greedy","leading_eigen","leiden","louvain","spinglass"] 
for com_detec_met in list_com_detec_met:

	#Column names the table
	com_col_name = "mem_{}".format(com_detec_met)	
	n_com = DFcom["mem_{}".format(com_detec_met)].max()

	for c in range(1,n_com+1):

		#Set the position of the figure
		num = (c-1)%3 #Number in a slide
		row = num//2 #Columns
		col = num%2 #Row

		#Make a new slide
		if num == 0:
			slide = prs.slides.add_slide(blank_slide_layout)
			title1 = "Network communities and their taxonomic composition"
			title2 = "{} - {} - {} - frq{} - {} - {}".format(SIZ,opo2,tpr2,FRQ,NRM,com_detec_met)
			tf = MakeTextBox(top=Inches(0.1), left=Inches(0.5), text=title1, size=Pt(30))
			AddTextToBox(tf, title2, size=Pt(24))

			#Add world map of communities
			img_path = "../figures/{}/V316.class.map.{}.{}.png".format(subdirec,NRM,com_detec_met)
			top_base = Inches(3.5*1 + 1.2)
			left_base = Inches(7.5*1 + 0.5)
			pic = slide.shapes.add_picture(img_path, left=left_base, top=top_base+Inches(0.5), height=Inches(3.5))
	
		#Add community name
		top_base = Inches(3.5*row + 1.2)
		left_base = Inches(7.5*col + 0.5)
		MakeTextBox(top=top_base, left=left_base, text="community {}".format(c), size=Pt(24))

		#Add image of subnetwork
		img_path = "../figures/{}/community/V314.subnetwork.{}.{}.community{}.png".format(subdirec,NRM,com_detec_met,c)
		pic = slide.shapes.add_picture(img_path, left=left_base, top=top_base+Inches(0.5), height=Inches(3.0))
	
		#Add image of taxonomic breakdown
		img_path = "../figures/{}/community/V314.taxonomy.breakdown.{}.{}.community{}.png".format(subdirec,NRM,com_detec_met,c)
		pic = slide.shapes.add_picture(img_path, left=left_base+Inches(3.5), top=top_base+Inches(0.5), height=Inches(3.0))

#Save the presentation
prs.save("../figures/{}/V318.summary.community.detail.{}.pptx".format(subdirec,NRM))
