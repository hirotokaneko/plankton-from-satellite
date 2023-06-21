import sys
from pptx import Presentation
from pptx.util import Inches, Pt

#Arguments
SIZ = sys.argv[1]
OPO = sys.argv[2]
TPR = sys.argv[3]
FRQ = sys.argv[4]
NRM = sys.argv[5]

#Set file name variables
opo2 = "" if OPO == "all" else "{}.".format(OPO)
tpr2 = "" if TPR == "keep" else "thinned{}.".format(TPR)
subdirec = "{}.{}{}frq{}".format(SIZ,opo2,tpr2,FRQ)

#Make a presentation
prs = Presentation()
prs.slide_height = Inches(9)
prs.slide_width = Inches(16)

#Set balank slide layout
blank_slide_layout = prs.slide_layouts[6]
font = "Helvetica"

#Read modularity data
with open("../data/{}/modularity.{}.{}.{}.{}.{}.txt".format(subdirec,SIZ,opo2,tpr2,FRQ,NRM)) as f:
	listModulariity = []
	gate = "closed"
	for line in f:
		val = line.strip()[4:]
		if val == '"Detect community"':
			gate = "open"
		if gate == "open":
			if val[0] != '"':
				listModulariity.append(float(val))

#Function for making a text box
def MakeTextBox(top, left, text, size):
	width = height = Inches(1)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = text
	p.font.size = size
	p.font.name = font
	return tf

#Function for adding a row to the text box
def AddTextToBox(tf, text, size):
	p = tf.add_paragraph()
	p.text = text
	p.font.size = size
	p.font.name = font

#Loop by community detection methods
list_com_detec_met = ["fast_greedy","infomap","label_prop","leading_eigen",
		      "leiden","louvain","spinglass","walktrap"] 
for i, com_detec_met in enumerate(list_com_detec_met):

	#Set the position of the figure
	num = i%4 #Number in a slide
	row = num//2 #Column
	col = num%2 #Row

	#Make a new slide
	if num == 0:
		slide = prs.slides.add_slide(blank_slide_layout)
		title1 = "Comparison of community detection algorithms"
		title2 = "{} - {} - {} - frq{} - {}".format(SIZ,opo2,tpr2,FRQ,NRM)
		tf = MakeTextBox(top=Inches(0.1), left=Inches(0.5), text=title1, size=Pt(30))
		AddTextToBox(tf, title2, size=Pt(24))
	
	#Add community detection method
	top_base = Inches(3.5*row + 1.2)
	left_base = Inches(7.5*col + 0.5)
	tf = MakeTextBox(top=top_base, left=left_base, text=com_detec_met, size=Pt(24))
	#Add modularity
	acc_text = "Modulariity = {0:.2f}".format(listModulariity[i])
	AddTextToBox(tf=tf, text=acc_text, size=Pt(20))
	
	#Add network image to the slide
	img_path = "../figures/{}/V313.network.{}.{}.png".format(subdirec,NRM,com_detec_met)
	pic = slide.shapes.add_picture(img_path, left=left_base, top=top_base+Inches(0.9), height=Inches(3.0))

	#Add edge statisfaction image to the slide
	img_path = "../figures/{}/V315.edge.satisfaction.heatmap.{}.{}.png".format(subdirec,NRM,com_detec_met)
	pic = slide.shapes.add_picture(img_path, left=left_base+Inches(4.5), top=top_base+Inches(0.9), height=Inches(2.5))

	#Add community breakdown image to the slide
	img_path = "../figures/{}/V313.community.breakdown.{}.{}.png".format(subdirec,NRM,com_detec_met)
	pic = slide.shapes.add_picture(img_path, left=left_base+Inches(5.5), top=top_base, height=Inches(1.2))

#Save the presentation
prs.save("../figures/{}/V317.summary.class.detection.{}.pptx".format(subdirec,NRM))
