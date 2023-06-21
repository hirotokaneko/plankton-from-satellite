import sys
from pptx import Presentation
from pptx.util import Inches, Pt

#Arguments
SIZ = sys.argv[1]
OPO = sys.argv[2]
TPR = sys.argv[3]
FRQ = sys.argv[4]
NRM = sys.argv[5]
COM = sys.argv[6]
MLM = sys.argv[7]
CTF = sys.argv[8]

#Set file name variables
opo2 = "" if OPO == "all" else "{}.".format(OPO)
tpr2 = "" if TPR == "keep" else "thinned{}.".format(TPR)
subdirec = "{}.{}{}frq{}".format(SIZ,opo2,tpr2,FRQ)

#Make a presentation
prs = Presentation()
prs.slide_height = Inches(9)
prs.slide_width = Inches(16)

#Set balank slide layout
blank_slide_layout = prs.slide_layouts[6]
font = "Helvetica"

#Parse the fit results and get the best estimators
list_feature_set = []
list_fit_model = []
flag = 0
with open("../data/{}/fit/fit.model.{}.{}.{}.txt".format(subdirec,NRM,COM,MLM)) as f:
	line_container = []
	for i, line in enumerate(f):
		line_container.append(line.strip())
		if "Best estimator" in line:
			list_feature_set.append(line_container[i-1])	
			if flag:
				texts = "".join(line_container[flag+2:i-2])
				list_fit_model.append(texts.split(",",1)[1].strip()[:-3])
			flag = i
	texts = "".join(line_container[flag+2:-1])
	list_fit_model.append(texts.split(",",1)[1].strip()[:-3])

print(list_feature_set)
print(list_fit_model)

#Function for making a text box
def MakeTextBox(top, left, text, size):
	width = height = Inches(1)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = text
	p.font.size = size
	p.font.name = font
	return tf

#Function for adding a row to the text box
def AddTextToBox(tf, text, size):
	p = tf.add_paragraph()
	p.text = text
	p.font.size = size
	p.font.name = font

#Loop by feature set
for i, fts in enumerate(list_feature_set):

	#Set the position of the figure
	row = i%3 #Columns

	#Make a new slide
	if row == 0:
		slide = prs.slides.add_slide(blank_slide_layout)
		title = "{} - {} - {} - frq{} - {} - {} - {}".format(SIZ,opo2,tpr2,FRQ,NRM,COM,MLM)
		MakeTextBox(top=Inches(0.1), left=Inches(0.5), text=title, size=Pt(30))
		MakeTextBox(top=Inches(0.7), left=Inches(3), text="Feature Importance", size=Pt(20))
		MakeTextBox(top=Inches(0.7), left=Inches(6.5), text="F1-score of CV", size=Pt(20))
		MakeTextBox(top=Inches(0.7), left=Inches(10), text="CV map when threshold > {}".format(CTF), size=Pt(20))

	#Set the top mergin
	top = Inches(2.5*row + 1.2)

	#Add the name of feature set
	tf = MakeTextBox(top=top, left=Inches(0.5), text=fts, size=Pt(24))
	#Add the best estimator
	for comma in list_fit_model[i].split(","):
		AddTextToBox(tf=tf, text=comma, size=Pt(16))
		
	#Add a image of feature importance
	img_path = "../figures/{}/fit/V412.feature.importance.{}.{}.{}.{}.png".format(subdirec,NRM,COM,MLM,fts)
	left = Inches(3)
	height = Inches(2.5)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)

	#Add a image of F-score
	img_path = "../figures/{}/cv/V411.Fscore.spatial0.{}.{}.{}.{}.png".format(subdirec,NRM,COM,MLM,fts)
	left = Inches(6.5)
	height = Inches(2.5)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)

	#Add a world map of predicted/not predicted
	img_path = "../figures/{}/cv/V413.prediction.map.spatial0.{}.{}.{}.{}.png".format(subdirec,NRM,COM,MLM,fts)
	left = Inches(10)
	height = Inches(2.5)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)

#Save the presentation
prs.save("../figures/{}/fit/V416.Summary.Fit.{}.{}.{}.pptx".format(subdirec,NRM,COM,MLM))

