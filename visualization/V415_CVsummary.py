import pandas as pd
import sys
from pptx import Presentation
from pptx.util import Inches, Pt

#Arguments
SIZ = sys.argv[1]
OPO = sys.argv[2]
TPR = sys.argv[3]
FRQ = sys.argv[4]
NRM = sys.argv[5]
COM = sys.argv[6]
MLM = sys.argv[7]
RAD = sys.argv[8]

#Set file name variables
opo2 = "" if OPO == "all" else "{}.".format(OPO)
tpr2 = "" if TPR == "keep" else "thinned{}.".format(TPR)
subdirec = "{}.{}{}frq{}".format(SIZ,opo2,tpr2,FRQ)

#Make a presentation
prs = Presentation()
prs.slide_height = Inches(9)
prs.slide_width = Inches(16)

#Set balank slide layout
blank_slide_layout = prs.slide_layouts[6]
font = "Helvetica"

#Read accuracy data
DFacc = pd.read_csv("../data/{}/cv/accuracy.spatial{}.{}.{}.{}.txt".format(subdirec,RAD,NRM,COM,MLM), index_col=0, delimiter="\t")
list_feature_set = list(DFacc.index)

#Function for making a text box
def MakeTextBox(top, left, text, size):
	width = height = Inches(1)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = text
	p.font.size = size
	p.font.name = font
	return tf

#Function for adding a row to the text box
def AddTextToBox(tf, text, size):
	p = tf.add_paragraph()
	p.text = text
	p.font.size = size
	p.font.name = font

#Loop by feature set
for i, fts in enumerate(list_feature_set):

	#Set the position of the figure
	row = i%3 #Columns

	#Make a new slide
	if row == 0:
		slide = prs.slides.add_slide(blank_slide_layout)
		title = "{} - {} - {} - frq{} - {} - {} - {} - spatial{}".format(SIZ,opo2,tpr2,FRQ,NRM,COM,MLM,RAD)
		MakeTextBox(top=Inches(0.1), left=Inches(0.5), text=title, size=Pt(30))
		MakeTextBox(top=Inches(0.7), left=Inches(3), text="Confusion Matrix", size=Pt(20))
		MakeTextBox(top=Inches(0.7), left=Inches(6.5), text="ROC Curve", size=Pt(20))
		MakeTextBox(top=Inches(0.7), left=Inches(10.7), text="Precision-Recall Curve", size=Pt(20))

	#Set the top mergin
	top = Inches(2.5*row + 1.2)

	#Add the name of feature set
	tf = MakeTextBox(top=top, left=Inches(0.5), text=fts, size=Pt(24))
	#Add accuracy value
	acc_text = "CV acc = {0:.2f}".format(DFacc.loc[fts,"accuracy"])
	AddTextToBox(tf=tf, text=acc_text, size=Pt(20))

	#Add cofusion matrix
	img_path = "../figures/{}/cv/V411.confusion.matrix.spatial{}.{}.{}.{}.{}.png".format(subdirec,RAD,NRM,COM,MLM,fts)
	left = Inches(3)
	height = Inches(2.5)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)

	#Add ROC curve
	img_path = "../figures/{}/cv/V411.ROCcurve.spatial{}.{}.{}.{}.{}.png".format(subdirec,RAD,NRM,COM,MLM,fts)
	left = Inches(6.5)
	height = Inches(2.5)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)

	#Add PR curve
	img_path = "../figures/{}/cv/V411.PRcurve.spatial{}.{}.{}.{}.{}.png".format(subdirec,RAD,NRM,COM,MLM,fts)
	left = Inches(10.7)
	height = Inches(2.5)
	pic = slide.shapes.add_picture(img_path, left, top, height=height)

#Save the presentation
prs.save("../figures/{}/cv/V415.summary.CV.spatial{}.{}.{}.{}.pptx".format(subdirec,RAD,NRM,COM,MLM))
