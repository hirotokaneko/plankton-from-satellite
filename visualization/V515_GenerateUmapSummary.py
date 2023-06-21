import sys
from pptx import Presentation
from pptx.util import Inches, Pt

#Arguments
SIZ = sys.argv[1]
OPO = sys.argv[2]
TPR = sys.argv[3]
FRQ = sys.argv[4]
AAJ = sys.argv[5]
NRM = sys.argv[6]

#Set file name variables
opo2 = "" if OPO == "all" else "{}.".format(OPO)
tpr2 = "" if TPR == "keep" else "thinned{}.".format(TPR)
subdirec = "{}.{}{}frq{}".format(SIZ,opo2,tpr2,FRQ)

#Make a presentation
prs = Presentation()
prs.slide_height = Inches(9)
prs.slide_width = Inches(16)

#Set balank slide layout
blank_slide_layout = prs.slide_layouts[6]
font = "Helvetica"

#Function for making a text box
def MakeTextBox(top, left, text, size):
	width = height = Inches(1)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.text = text
	p.font.size = size
	p.font.name = font
	return tf

#Function for adding a row to the text box
def AddTextToBox(tf, text, size):
	p = tf.add_paragraph()
	p.text = text
	p.font.size = size
	p.font.name = font

#Loop by community detection methods
list_com_detec_met = ["fast_greedy","leading_eigen","leiden","louvain","spinglass"] 
for i, com_detec_met in enumerate(list_com_detec_met):

	#Make a new slide
	slide = prs.slides.add_slide(blank_slide_layout)
	title1 = "Assigned class on U-map and SST"
	title2 = "{} - {} - {} - frq{} - {} - {}".format(SIZ,opo2,tpr2,FRQ,NRM,com_detec_met)
	tf = MakeTextBox(top=Inches(0.1), left=Inches(0.5), text=title1, size=Pt(30))
	AddTextToBox(tf, title2, size=Pt(24))

	#Loop by feature sets
	feature_sets = ["space","product","rrs","satellite","allfeat"]
	for i, fts in enumerate(feature_sets):

		#Set the position of the figure
		row = i//3 #Columns
		col = i%3 #Row
	
		#Add the name of feature set
		top_base = Inches(3.5*row + 1.2)
		left_base = Inches(5*col + 0.5)
		MakeTextBox(top=top_base, left=left_base, text=fts, size=Pt(24))
	
		#Add a image of U-map
		img_path = "../figures/{}/V513.class.umap.{}.{}.{}.{}.png".format(subdirec,AAJ,NRM,com_detec_met,fts)
		pic = slide.shapes.add_picture(img_path, left=left_base, top=top_base+Inches(0.5), height=Inches(3.0))

	#Add "sst"
	top_base = Inches(3.5*1 + 1.2)
	left_base = Inches(5*2 + 0.5)
	MakeTextBox(top=top_base, left=left_base, text="sst", size=Pt(24))
	#Add aimage of "sst"
	img_path = "../figures/{}/V514.class.sst.{}.{}.png".format(subdirec,NRM,com_detec_met)
	pic = slide.shapes.add_picture(img_path, left=left_base, top=top_base+Inches(0.5), height=Inches(3.0))

#Save the presentation
prs.save("../figures/{}/V515.summary.umap.{}.{}.pptx".format(subdirec,AAJ,NRM))
